import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_de_presentation():
    prs = Presentation()
    # 設定 16:9 寬螢幕比例
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 顏色配置 (Tech Dark & Clean Accent)
    DARK_BG = RGBColor(30, 35, 45)      # 深灰色背景
    TEXT_WHITE = RGBColor(255, 255, 255)
    TEXT_LIGHT_GRAY = RGBColor(180, 185, 195)
    ACCENT_GREEN = RGBColor(46, 204, 113) # dbt 亮綠色

    # 資料定義
    slides_data = [
        {
            "title": "現代數據棧（MDS）實戰\n端到端巴西電商數據管線與自動化數據防禦系統",
            "subtitle": "報告人：資料工程師 | 利用 Python, PostgreSQL 與 dbt 打造生產級架構",
            "type": "title"
        },
        {
            "title": "數據背景與企業痛點分析",
            "points": [
                "真實大數據場景：採用 Olist 巴西電商真實數據集，涵蓋訂單、客戶、商品明細等多源高關聯性資料（十萬級別數據量）。",
                "傳統 ETL 的效能瓶頸：原始資料欄位繁雜，傳統思維若盲目使用 SELECT * 進行多表 JOIN，會導致大量的內存（Memory）浪費、嚴重拖慢資料庫執行計畫（Execution Plan）效能。",
                "商業痛點：缺乏統一的資料維度與主資料（Master Data）管理，導致後端分析師在計算指標時必須撰寫冗長的 SQL，且資料品質（如空值、重複值）無法得到有效監控。"
            ],
            "type": "content"
        },
        {
            "title": "端到端現代數據棧（MDS）架構設計",
            "points": [
                "高速 Ingestion (EL 層)：使用 Python 自動化腳本，將多張原始 CSV 數據高效導入 PostgreSQL 18 生產環境。",
                "兩層流俐落架構 (Clean Code)：拒絕教科書式的三層冗餘 CTE。在 Staging 層（清洗層）直接精簡為兩層流結構，畫面乾淨、易於維護。",
                "精準欄位過濾（脂肪退散）：堅決捨棄 SELECT * 的壞習慣。在 Staging 層即實施「抓大放小」策略，僅提取核心分析欄位，從源頭優化查詢效率。"
            ],
            "note": "💡 [提示：請在右側預留空間貼上你的 dbt Lineage DAG 綠色網頁截圖]",
            "type": "content"
        },
        {
            "title": "以 Kimball 理論構築電商星狀模型 (Star Schema)",
            "points": [
                "主資料特性管理 (Dimension Table)：提煉 stg_customers，建立以 customer_unique_id 為核心的客戶主資料維度表：dim_customers。完美對齊傳統資料倉儲的主資料屬性管理思維。",
                "高性能事實表 (Fact Table)：打造核心商業事實大表：fct_orders。精準橫跨訂單層、明細層與客戶層，並將其實體化（Materialized as Table），確保前端 BI 查詢高頻調用時的速度。",
                "核心商業指標提煉：在事實表中直接衍生計算 total_order_amount（商品單價 + 運費），將複雜的商業邏輯封裝在數據倉庫內部，向下游提供乾淨、單一事實來源（Single Source of Truth）。"
            ],
            "type": "content"
        },
        {
            "title": "宣告式測試與數據品質（Data Quality）防線",
            "points": [
                "告別人工對帳：引入 dbt 宣告式測試（Declarative Testing）機制，將數據品質校驗（Data Quality）全面自動化。",
                "主鍵唯一性防禦 (PK Validation)：針對維度表的主鍵進行 unique 與 not_null 測試，確保主資料庫無重複與遺漏。",
                "商業邊界測試 (Boundary Value Test)：針對訂單狀態（order_status）導入 accepted_values 測試，嚴格限制非預期狀態資料流入下游，保證 BI 層數據 100% 準確。"
            ],
            "note": "💡 [提示：請在右側預留空間貼上你的 dbt test 一整排綠色 PASS 的終端機截圖]",
            "type": "content"
        },
        {
            "title": "數據價值呈現：商業決策大盤 (BI Dashboard)",
            "points": [
                "數據驅動商業決策：前端看板直接對接 PostgreSQL 裡的 fct_orders，實現地域營收分析、平均客單價追蹤、以及各城市運費成本佔比分析。",
                "完備的數據生命週期：本專案成功實現了從 原始數據 -> 清洗建模 -> 品質防禦 -> 商業價值可視化 的完整閉環，體現現代數據工程師（Modern Data Engineer）的全面技術素養與業務洞察力。"
            ],
            "note": "💡 [提示：請在中央/右側空間貼上我們接下來要拉的 Power BI 看板截圖]",
            "type": "content"
        }
    ]

    for data in slides_data:
        # 使用空白版面，完全靠程式碼控制排版
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # 填滿深色背景
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG

        if data["type"] == "title":
            # 封面大標題
            txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(2))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = data["title"]
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = TEXT_WHITE
            p.font.name = 'Microsoft JhengHei'

            # 封面副標題
            p2 = tf.add_paragraph()
            p2.text = "\n" + data["subtitle"]
            p2.font.size = Pt(18)
            p2.font.color.rgb = ACCENT_GREEN
            p2.font.name = 'Microsoft JhengHei'

        else:
            # 內頁標題
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.5), Inches(1))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = data["title"]
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = ACCENT_GREEN
            p.font.name = 'Microsoft JhengHei'

            # 內頁內容 (左側)
            contentBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(7.5), Inches(5))
            ctf = contentBox.text_frame
            ctf.word_wrap = True

            for i, pt in enumerate(data["points"]):
                cp = ctf.paragraphs[0] if i == 0 else ctf.add_paragraph()
                cp.text = "• " + pt
                cp.font.size = Pt(15)
                cp.font.color.rgb = TEXT_WHITE
                cp.font.name = 'Microsoft JhengHei'
                cp.space_after = Pt(14)

            # 如果有備註提示 (放右下角)
            if "note" in data:
                noteBox = slide.shapes.add_textbox(Inches(8.5), Inches(5.5), Inches(4), Inches(1))
                ntf = noteBox.text_frame
                ntf.word_wrap = True
                np = ntf.paragraphs[0]
                np.text = data["note"]
                np.font.size = Pt(12)
                np.font.color.rgb = TEXT_LIGHT_GRAY
                np.font.name = 'Microsoft JhengHei'

    output_filename = "modern_de_project_presentation.pptx"
    prs.save(output_filename)
    print(f"🎉 簡報成功誕生！檔案儲存為: {os.path.abspath(output_filename)}")

if __name__ == "__main__":
    create_de_presentation()